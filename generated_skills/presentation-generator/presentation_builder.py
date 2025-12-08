"""
Presentation Builder Module.
Generates PowerPoint presentations from structured content.
Requires: pip install python-pptx Pillow requests
"""

from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass, field
from enum import Enum
import re
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("python-pptx not installed. Run: pip install python-pptx")


class SlideType(Enum):
    """Types of slides."""
    TITLE = "title"
    SECTION = "section"
    CONTENT = "content"
    TWO_COLUMN = "two_column"
    IMAGE = "image"
    QUOTE = "quote"
    DATA = "data"
    COMPARISON = "comparison"
    BLANK = "blank"


class PresentationStyle(Enum):
    """Visual styles for presentations."""
    CORPORATE = "corporate"
    CREATIVE = "creative"
    MINIMAL = "minimal"
    BOLD = "bold"
    DARK = "dark"


@dataclass
class SlideContent:
    """Content for a single slide."""
    slide_number: int
    slide_type: SlideType
    headline: str
    body: List[str] = field(default_factory=list)
    speaker_notes: str = ""
    image_prompt: str = ""
    stock_keywords: str = ""
    chart_description: str = ""
    subtitle: str = ""


@dataclass
class PresentationTheme:
    """Visual theme settings."""
    primary_color: Tuple[int, int, int] = (0, 102, 204)  # Blue
    secondary_color: Tuple[int, int, int] = (0, 51, 102)  # Navy
    accent_color: Tuple[int, int, int] = (40, 167, 69)  # Green
    background_color: Tuple[int, int, int] = (255, 255, 255)  # White
    text_color: Tuple[int, int, int] = (51, 51, 51)  # Dark gray
    font_title: str = "Segoe UI"
    font_body: str = "Segoe UI"
    title_size: int = 44
    subtitle_size: int = 28
    body_size: int = 18


# Predefined themes
THEMES = {
    PresentationStyle.CORPORATE: PresentationTheme(
        primary_color=(0, 102, 204),
        secondary_color=(0, 51, 102),
        text_color=(51, 51, 51),
    ),
    PresentationStyle.CREATIVE: PresentationTheme(
        primary_color=(156, 39, 176),
        secondary_color=(103, 58, 183),
        accent_color=(255, 152, 0),
        text_color=(33, 33, 33),
    ),
    PresentationStyle.MINIMAL: PresentationTheme(
        primary_color=(33, 33, 33),
        secondary_color=(97, 97, 97),
        accent_color=(0, 150, 136),
        text_color=(66, 66, 66),
    ),
    PresentationStyle.BOLD: PresentationTheme(
        primary_color=(244, 67, 54),
        secondary_color=(33, 33, 33),
        accent_color=(255, 235, 59),
        text_color=(33, 33, 33),
        title_size=54,
    ),
    PresentationStyle.DARK: PresentationTheme(
        primary_color=(66, 165, 245),
        secondary_color=(41, 121, 255),
        background_color=(33, 33, 33),
        text_color=(255, 255, 255),
    ),
}


class PresentationBuilder:
    """Build PowerPoint presentations programmatically."""

    def __init__(self, style: PresentationStyle = PresentationStyle.CORPORATE):
        """
        Initialize presentation builder.

        Args:
            style: Visual style to apply
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required. Install with: pip install python-pptx")

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)
        self.style = style
        self.theme = THEMES.get(style, THEMES[PresentationStyle.CORPORATE])
        self.slides: List[SlideContent] = []

    def _rgb(self, color: Tuple[int, int, int]) -> RgbColor:
        """Convert RGB tuple to RgbColor."""
        return RgbColor(color[0], color[1], color[2])

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        author: str = "",
        date: str = "",
    ) -> None:
        """Add title slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Background shape
        left = Inches(0)
        top = Inches(0)
        width = self.prs.slide_width
        height = self.prs.slide_height

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(self.theme.primary_color)
        shape.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(self.theme.title_size)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor(255, 255, 255)
        title_para.font.name = self.theme.font_title
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
            )
            sub_frame = sub_box.text_frame
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = subtitle
            sub_para.font.size = Pt(self.theme.subtitle_size)
            sub_para.font.color.rgb = RgbColor(255, 255, 255)
            sub_para.font.name = self.theme.font_body
            sub_para.alignment = PP_ALIGN.CENTER

        # Author and date
        if author or date:
            info_text = f"{author}  |  {date}" if author and date else author or date
            info_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5)
            )
            info_frame = info_box.text_frame
            info_para = info_frame.paragraphs[0]
            info_para.text = info_text
            info_para.font.size = Pt(14)
            info_para.font.color.rgb = RgbColor(200, 200, 200)
            info_para.font.name = self.theme.font_body
            info_para.alignment = PP_ALIGN.CENTER

    def add_section_slide(self, title: str, subtitle: str = "") -> None:
        """Add section divider slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(3), Inches(13.333), Inches(1.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(self.theme.primary_color)
        bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(3.2), Inches(11.833), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor(255, 255, 255)
        title_para.font.name = self.theme.font_title

    def add_content_slide(
        self,
        headline: str,
        bullets: List[str],
        speaker_notes: str = "",
    ) -> None:
        """Add content slide with bullet points."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Header bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(self.theme.primary_color)
        bar.line.fill.background()

        # Headline
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = headline
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor(255, 255, 255)
        title_para.font.name = self.theme.font_title

        # Bullets
        bullet_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.8), Inches(11.833), Inches(5)
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(self.theme.body_size)
            p.font.color.rgb = self._rgb(self.theme.text_color)
            p.font.name = self.theme.font_body
            p.space_after = Pt(12)

        # Speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def add_quote_slide(self, quote: str, author: str = "") -> None:
        """Add quote slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Quote marks decorative
        quote_mark = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(2), Inches(2)
        )
        qm_frame = quote_mark.text_frame
        qm_para = qm_frame.paragraphs[0]
        qm_para.text = """
        qm_para.font.size = Pt(120)
        qm_para.font.color.rgb = self._rgb(self.theme.primary_color)

        # Quote text
        quote_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(3)
        )
        q_frame = quote_box.text_frame
        q_para = q_frame.paragraphs[0]
        q_para.text = quote
        q_para.font.size = Pt(28)
        q_para.font.italic = True
        q_para.font.color.rgb = self._rgb(self.theme.text_color)
        q_para.font.name = self.theme.font_body

        # Author
        if author:
            auth_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.5), Inches(11.333), Inches(0.5)
            )
            a_frame = auth_box.text_frame
            a_para = a_frame.paragraphs[0]
            a_para.text = f"— {author}"
            a_para.font.size = Pt(18)
            a_para.font.color.rgb = self._rgb(self.theme.secondary_color)
            a_para.alignment = PP_ALIGN.RIGHT

    def add_two_column_slide(
        self,
        headline: str,
        left_content: List[str],
        right_content: List[str],
        left_title: str = "",
        right_title: str = "",
    ) -> None:
        """Add two-column comparison slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Header bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(self.theme.primary_color)
        bar.line.fill.background()

        # Headline
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = headline
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor(255, 255, 255)

        # Left column
        self._add_column(slide, Inches(0.5), left_title, left_content)

        # Right column
        self._add_column(slide, Inches(6.9), right_title, right_content)

    def _add_column(
        self, slide, left_pos, title: str, content: List[str]
    ) -> None:
        """Add a column to a slide."""
        if title:
            title_box = slide.shapes.add_textbox(
                left_pos, Inches(1.5), Inches(5.8), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self._rgb(self.theme.primary_color)

        content_box = slide.shapes.add_textbox(
            left_pos, Inches(2.2), Inches(5.8), Inches(4.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = self._rgb(self.theme.text_color)
            p.space_after = Pt(8)

    def add_closing_slide(
        self,
        title: str = "Thank You",
        subtitle: str = "",
        contact: str = "",
    ) -> None:
        """Add closing/thank you slide."""
        self.add_title_slide(title, subtitle, contact)

    def save(self, filepath: str) -> str:
        """
        Save presentation to file.

        Args:
            filepath: Output file path (.pptx)

        Returns:
            Absolute path to saved file
        """
        if not filepath.endswith('.pptx'):
            filepath += '.pptx'

        self.prs.save(filepath)
        return os.path.abspath(filepath)


def generate_image_prompts(topic: str, slide_titles: List[str]) -> Dict[str, str]:
    """
    Generate AI image prompts for each slide.

    Args:
        topic: Main presentation topic
        slide_titles: List of slide headlines

    Returns:
        Dictionary mapping slide titles to image prompts
    """
    prompts = {}

    base_style = "professional, corporate style, clean background, high quality, 16:9 aspect ratio"

    for title in slide_titles:
        # Create contextual prompt based on title keywords
        title_lower = title.lower()

        if any(word in title_lower for word in ["intro", "title", "welcome"]):
            prompt = f"Abstract professional background representing {topic}, {base_style}, subtle gradients"
        elif any(word in title_lower for word in ["team", "people", "who"]):
            prompt = f"Diverse professional team collaborating in modern office, {base_style}"
        elif any(word in title_lower for word in ["data", "chart", "number", "result"]):
            prompt = f"Abstract data visualization, flowing charts and graphs, {base_style}, blue tones"
        elif any(word in title_lower for word in ["future", "vision", "next"]):
            prompt = f"Futuristic concept representing innovation and progress, {base_style}, inspiring"
        elif any(word in title_lower for word in ["problem", "challenge", "issue"]):
            prompt = f"Visual metaphor of overcoming obstacles, {base_style}, dramatic lighting"
        elif any(word in title_lower for word in ["solution", "answer", "how"]):
            prompt = f"Light bulb or puzzle pieces coming together, {base_style}, success imagery"
        elif any(word in title_lower for word in ["thank", "question", "contact"]):
            prompt = f"Professional handshake or thank you gesture, {base_style}, warm tones"
        else:
            prompt = f"Professional visual representing {title}, {base_style}, relevant to {topic}"

        prompts[title] = prompt

    return prompts


# Example usage
if __name__ == "__main__":
    if PPTX_AVAILABLE:
        # Create presentation
        builder = PresentationBuilder(style=PresentationStyle.CORPORATE)

        # Add slides
        builder.add_title_slide(
            title="Introduction to AI",
            subtitle="A Guide for Business Leaders",
            author="Your Company",
            date="November 2025"
        )

        builder.add_section_slide("Why AI Matters")

        builder.add_content_slide(
            headline="AI is Transforming Every Industry",
            bullets=[
                "Healthcare: 40% faster diagnosis with AI imaging",
                "Finance: Fraud detection accuracy improved by 60%",
                "Retail: Personalization increases conversion by 35%",
                "Manufacturing: Predictive maintenance reduces downtime 25%"
            ],
            speaker_notes="Emphasize that AI is not future tech - it's here now and delivering real results."
        )

        builder.add_two_column_slide(
            headline="AI vs Traditional Approaches",
            left_title="Traditional",
            left_content=["Manual processes", "Reactive decisions", "Limited scale", "High error rate"],
            right_title="AI-Powered",
            right_content=["Automated workflows", "Predictive insights", "Unlimited scale", "Consistent accuracy"]
        )

        builder.add_quote_slide(
            quote="AI is probably the most important thing humanity has ever worked on.",
            author="Sundar Pichai, CEO Google"
        )

        builder.add_closing_slide(
            title="Questions?",
            subtitle="Let's discuss how AI can help your business",
            contact="email@company.com"
        )

        # Save
        output_path = builder.save("sample_presentation.pptx")
        print(f"Presentation saved to: {output_path}")

        # Generate image prompts
        slide_titles = [
            "Introduction to AI",
            "Why AI Matters",
            "AI is Transforming Every Industry",
            "AI vs Traditional Approaches",
            "Questions?"
        ]
        prompts = generate_image_prompts("Artificial Intelligence for Business", slide_titles)

        print("\nImage Prompts for DALL-E/Midjourney:")
        for title, prompt in prompts.items():
            print(f"\n{title}:")
            print(f"  {prompt}")
    else:
        print("Install python-pptx to run this example: pip install python-pptx")
